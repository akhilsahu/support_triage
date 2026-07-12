"""
PptxParser — python-pptx. Each slide → one ParsedPage.
Speaker notes are appended to the slide text.
"""

from __future__ import annotations

import io
import subprocess
import tempfile
from pathlib import Path
from typing import List

import structlog

from app.rag.document_parser import ParsedDocument, ParsedPage
from app.orchestra.ai.ingestion.core.base import BaseParser

logger = structlog.get_logger()


class PptxParser(BaseParser):
    extensions = [".pptx", ".ppt"]

    def parse(self, raw: bytes, filename: str) -> ParsedDocument:
        ext = Path(filename).suffix.lower()
        if ext == ".ppt":
            raw, filename = self._convert_ppt(raw, filename)

        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

        prs   = Presentation(io.BytesIO(raw))
        pages: List[ParsedPage] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            parts: List[str] = []

            # Slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(f"Slide {slide_num}: {slide.shapes.title.text.strip()}")

            # All text shapes (tables + text boxes)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                        parts.append(text)

                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        parts.append("Table:\n" + "\n".join(rows))

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"Notes: {notes}")

            if parts:
                section = slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {slide_num}"
                pages.append(ParsedPage(
                    page=slide_num,
                    text="\n\n".join(parts),
                    section=section,
                ))

        meta = {"slide_count": len(prs.slides)}
        logger.info("ingestion.pptx", filename=filename, slides=len(pages))
        return ParsedDocument(filename=filename, extension=".pptx", pages=pages, metadata=meta)

    def _convert_ppt(self, raw: bytes, filename: str):
        if not self.cfg.libreoffice_enabled:
            raise ValueError(".ppt files require LibreOffice. Set INGESTION_LIBREOFFICE_ENABLED=true")
        with tempfile.TemporaryDirectory() as tmpdir:
            src = Path(tmpdir) / filename
            src.write_bytes(raw)
            try:
                subprocess.run(
                    [self.cfg.libreoffice_path, "--headless", "--convert-to", "pptx",
                     "--outdir", tmpdir, str(src)],
                    check=True, capture_output=True, timeout=60,
                )
            except (subprocess.CalledProcessError, FileNotFoundError) as e:
                raise RuntimeError(f"LibreOffice conversion failed for '{filename}': {e}")
            out = src.with_suffix(".pptx")
            if not out.exists():
                raise RuntimeError(f"LibreOffice did not produce output for '{filename}'")
            return out.read_bytes(), out.name
